"""
Generate BAREAI pre-panel PowerPoint (.pptx) with branded design.
Output: pre-panel/BAREAI_PrePanel_Presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "pre-panel"
SCREENSHOTS = ROOT / "thesis_assets" / "screenshots"
OUT_FILE = OUT_DIR / "BAREAI_PrePanel_Presentation.pptx"

# BAREAI brand palette
TEAL = RGBColor(0x0F, 0x76, 0x6E)
TEAL_DARK = RGBColor(0x0D, 0x5C, 0x56)
SLATE = RGBColor(0x0F, 0x17, 0x2A)
SLATE_LIGHT = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)

AUTHORS = (
    "Yaasiin Mohamuud Abdullaahi  |  Naima Abdiaziz Said\n"
    "Nasteha Mohamuud Mohamed  |  Najma Muhidiin Mohamed"
)
TITLE_EN = "Automatic Classification of Crime-Related Text Reports Using NLP"
SUBTITLE = "BAREAI — Bare AI Crime Intelligence Platform"


def inches(x: float) -> Inches:
    return Inches(x)


def add_bg(slide, color: RGBColor = SLATE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title: str, slide_num: int | None = None) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, inches(13.33), inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, inches(1.05), inches(0.12), inches(6.45))
    accent.fill.solid()
    accent.fill.fore_color.rgb = AMBER
    accent.line.fill.background()

    tb = slide.shapes.add_textbox(inches(0.45), inches(0.18), inches(10.5), inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"

    logo = slide.shapes.add_textbox(inches(11.2), inches(0.2), inches(1.8), inches(0.6))
    lp = logo.text_frame.paragraphs[0]
    lp.text = "BAREAI"
    lp.font.size = Pt(18)
    lp.font.bold = True
    lp.font.color.rgb = AMBER
    lp.alignment = PP_ALIGN.RIGHT

    if slide_num is not None:
        num = slide.shapes.add_textbox(inches(12.3), inches(6.9), inches(0.8), inches(0.35))
        np_ = num.text_frame.paragraphs[0]
        np_.text = str(slide_num)
        np_.font.size = Pt(11)
        np_.font.color.rgb = GRAY
        np_.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, items: list[str], left: float = 0.55, top: float = 1.35,
                width: float = 12.0, height: float = 5.5, size: int = 20) -> None:
    box = slide.shapes.add_textbox(inches(left), inches(top), inches(width), inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Segoe UI"
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
        p.line_spacing = 1.25


def add_note_box(slide, text: str, top: float = 5.8) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(0.55), inches(top), inches(12.2), inches(0.85))
    box.fill.solid()
    box.fill.fore_color.rgb = SLATE_LIGHT
    box.line.color.rgb = TEAL
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = AMBER
    p.font.name = "Segoe UI"


def add_image_if_exists(slide, filename: str, left: float, top: float, width: float) -> bool:
    path = SCREENSHOTS / filename
    if path.exists():
        slide.shapes.add_picture(str(path), inches(left), inches(top), width=inches(width))
        return True
    return False


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # decorative circles
    for x, y, sz, alpha in [(10.5, 0.3, 2.2, TEAL_DARK), (11.2, 4.5, 3.0, SLATE_LIGHT)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, inches(x), inches(y), inches(sz), inches(sz))
        c.fill.solid()
        c.fill.fore_color.rgb = alpha
        c.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, inches(13.33), inches(0.18))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TEAL
    top_bar.line.fill.background()

    uni = slide.shapes.add_textbox(inches(0.6), inches(0.55), inches(12), inches(0.5))
    up = uni.text_frame.paragraphs[0]
    up.text = "JAMHURIYA UNIVERSITY OF SCIENCE AND TECHNOLOGY (JUST)"
    up.font.size = Pt(14)
    up.font.color.rgb = GRAY
    up.font.name = "Segoe UI"
    up.alignment = PP_ALIGN.CENTER

    fac = slide.shapes.add_textbox(inches(0.6), inches(0.95), inches(12), inches(0.4))
    fp = fac.text_frame.paragraphs[0]
    fp.text = "Faculty of Computer & Information Technology"
    fp.font.size = Pt(13)
    fp.font.color.rgb = GRAY
    fp.font.name = "Segoe UI"
    fp.alignment = PP_ALIGN.CENTER

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(4.8), inches(1.55), inches(3.7), inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = TEAL
    badge.line.fill.background()
    bt = badge.text_frame.paragraphs[0]
    bt.text = "PRE-PANEL PRESENTATION"
    bt.font.size = Pt(14)
    bt.font.bold = True
    bt.font.color.rgb = WHITE
    bt.alignment = PP_ALIGN.CENTER

    brand = slide.shapes.add_textbox(inches(0.6), inches(2.35), inches(12), inches(0.9))
    bp = brand.text_frame.paragraphs[0]
    bp.text = "BAREAI"
    bp.font.size = Pt(54)
    bp.font.bold = True
    bp.font.color.rgb = TEAL
    bp.alignment = PP_ALIGN.CENTER

    tbox = slide.shapes.add_textbox(inches(0.8), inches(3.2), inches(11.7), inches(1.2))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = TITLE_EN
    tp.font.size = Pt(22)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER
    tp.font.name = "Segoe UI"

    sub = slide.shapes.add_textbox(inches(0.8), inches(4.35), inches(11.7), inches(0.5))
    sp = sub.text_frame.paragraphs[0]
    sp.text = SUBTITLE
    sp.font.size = Pt(16)
    sp.font.color.rgb = AMBER
    sp.alignment = PP_ALIGN.CENTER

    auth = slide.shapes.add_textbox(inches(1.0), inches(5.2), inches(11.3), inches(0.9))
    ap = auth.text_frame.paragraphs[0]
    ap.text = AUTHORS.replace("\n", "  •  ")
    ap.font.size = Pt(11)
    ap.font.color.rgb = GRAY
    ap.alignment = PP_ALIGN.CENTER

    date = slide.shapes.add_textbox(inches(0.8), inches(6.15), inches(11.7), inches(0.4))
    dp = date.text_frame.paragraphs[0]
    dp.text = "August 2026"
    dp.font.size = Pt(14)
    dp.font.bold = True
    dp.font.color.rgb = WHITE
    dp.alignment = PP_ALIGN.CENTER


def slide_content(prs: Presentation, num: int, title: str, bullets: list[str],
                  note: str | None = None, image: tuple[str, float, float, float] | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, title, num)

    if image and add_image_if_exists(slide, image[0], image[1], image[2], image[3]):
        add_bullets(slide, bullets, left=0.55, top=1.35, width=5.8, height=4.8, size=17)
    else:
        add_bullets(slide, bullets, size=20)

    if note:
        add_note_box(slide, note)


def slide_two_column(prs: Presentation, num: int, title: str,
                     left_title: str, left_items: list[str],
                     right_title: str, right_items: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, title, num)

    for i, (hdr, items, x) in enumerate([(left_title, left_items, 0.55), (right_title, right_items, 6.85)]):
        hbox = slide.shapes.add_textbox(inches(x), inches(1.25), inches(5.8), inches(0.45))
        hp = hbox.text_frame.paragraphs[0]
        hp.text = hdr
        hp.font.size = Pt(18)
        hp.font.bold = True
        hp.font.color.rgb = AMBER
        hp.font.name = "Segoe UI"

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(x), inches(1.75), inches(5.9), inches(4.9))
        card.fill.solid()
        card.fill.fore_color.rgb = SLATE_LIGHT
        card.line.color.rgb = TEAL

        add_bullets(slide, items, left=x + 0.2, top=1.95, width=5.5, height=4.5, size=16)


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "Qaab-dhismeedka Nidaamka (Architecture)", 6)

    layers = [
        ("PRESENTATION LAYER", "React 19 + Vite  |  Port 5173\nDashboard • Analysis • Cases • Reports", TEAL),
        ("APPLICATION LAYER", "Node.js + Express 5  |  Port 5000\nAuth • API • Facebook Monitor • Blacklist", RGBColor(0x1E, 0x40, 0xAF)),
        ("AI SERVICE", "Python Flask  |  Port 5001\ncrime_model.pkl • vectorizer.pkl", RGBColor(0x7C, 0x3A, 0xED)),
        ("DATA LAYER", "MongoDB — crime_detection_system\nUsers • Analyses • Cases • Notifications", RGBColor(0xB4, 0x53, 0x09)),
    ]
    y = 1.35
    for name, desc, color in layers:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(1.8), inches(y), inches(9.7), inches(1.15))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = name
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = WHITE
        y += 1.35

    add_note_box(slide, "JWT + bcrypt + RBAC  |  REST API  |  3-tier + AI microservice", top=6.0)


def slide_thank_you(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, TEAL_DARK)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, inches(3.0), inches(13.33), inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    t = slide.shapes.add_textbox(inches(0.5), inches(3.15), inches(12.3), inches(1.2))
    tp = t.text_frame.paragraphs[0]
    tp.text = "Mahadsanid!"
    tp.font.size = Pt(48)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER

    s = slide.shapes.add_textbox(inches(0.5), inches(4.75), inches(12.3), inches(0.6))
    sp = s.text_frame.paragraphs[0]
    sp.text = "Su'aalo?"
    sp.font.size = Pt(28)
    sp.font.color.rgb = AMBER
    sp.alignment = PP_ALIGN.CENTER

    c = slide.shapes.add_textbox(inches(0.5), inches(5.6), inches(12.3), inches(0.8))
    cp = c.text_frame.paragraphs[0]
    cp.text = "BAREAI — Bare AI Crime Intelligence Platform"
    cp.font.size = Pt(14)
    cp.font.color.rgb = GRAY
    cp.alignment = PP_ALIGN.CENTER


def build_pptx() -> Path:
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = inches(13.33)
    prs.slide_height = inches(7.5)

    slide_title(prs)

    slide_content(prs, 2, "Asalka Mowduuca (Background)", [
        "Dembiyada maalin walba way dhacaan — warbixinnadu waxay ku yimaadaan qoraal.",
        "NLP wuxuu kombiyuutarka u suurtageliyaa inuu fahmo luqadda dadka.",
        "Machine Learning wuxuu bartaa qaabka qoraallada dembiga iyo kuwa aan dembiga ahayn.",
        "Soomaaliya: qoraallo isku dhafan (Soomaali + Ingiriisi) — gacanta ma dabooli karto.",
    ], note="Hadaf: Nidaam casri ah oo ka caawiya hay'adaha amniga kala saarid degdeg ah.")

    slide_content(prs, 3, "Dhibaatada (Problem Statement)", [
        "Kala saarid GACANTA ah → gaabis, khaladaad, qaalin.",
        "Tusaale: 500 qoraal/maalin × 2 daq = 16+ saacadood kala saarid.",
        "Nidaamyo jira: Ingiriisi kaliya, keyword filter, ma jiro workflow.",
        "GAP: Ma jiro platform Soomaali + ML + monitoring + cases.",
    ], note="Su'aasha cilmi-baarista: Sidee loo dhiso nidaam NLP ah oo ku habboon Soomaaliya?")

    slide_two_column(prs, 4, "Ujeeddooyinka (Objectives)",
                     "Ujeeddo Guud", [
                         "In la dhiso framework otomaatig ah oo ML + NLP ku kala saara warbixinnada dembiga.",
                     ],
                     "Ujeeddooyin Gaar ah", [
                         "Yareynta ku tiirsanaanta gacanta",
                         "Uruurinta & diyaarinta xogta",
                         "Soo saarista features sax ah",
                         "Dhisidda nidaam NLP otomaatig ah",
                     ])

    slide_content(prs, 5, "Xalka: BAREAI", [
        "Platform web ah oo dhammaystiran — ma aha kaliya model.",
        "Kala saarid: Crime vs Not-crime + confidence score.",
        "Input: Text • URL • File (PDF/DOCX) • Batch.",
        "Hybrid: ML + ereyo Soomaali + goobo + blacklist.",
        "Facebook monitoring + Cases + Notifications + Reports.",
    ], image=("fig_5_16_result.png", 6.9, 1.35, 5.8))

    slide_architecture(prs)

    slide_content(prs, 7, "Habka Cilmi-baarista (Methodology)", [
        "Dataset: dataset.csv.csv — crime / not crime labels.",
        "Preprocessing: nadiifin → TF-IDF → train/test 80/20.",
        "Models: SVM • Random Forest • BERT (isbarbardhig).",
        "Metrics: Accuracy, Precision, Recall, F1, AUC-ROC.",
        "Kan ugu fiican → Flask /predict endpoint.",
    ])

    slide_content(prs, 8, "Astaamaha Muhiimka ah (Key Features)", [
        "Analysis Center — 4 nooc oo input ah.",
        "Dashboard — KPIs, trends, top keywords.",
        "Blacklist — pages, persons, keywords, websites.",
        "Notifications — digniino degdeg ah.",
        "Case Management — investigator + verdict.",
        "PDF Reports — weekly, monthly, individual.",
    ], image=("fig_5_15_analysis.png", 6.9, 1.35, 5.8))

    slide_two_column(prs, 9, "Natiijooyinka & Tijaabada",
                     "Model & Performance", [
                         "SVM, RF, BERT — isbarbardhig",
                         "Jawaabta AI: < 3 ilbiriqsi",
                         "End-to-end: < 30 ilbiriqsi",
                         "Keyword fallback haddii AI dhaco",
                     ],
                     "Testing", [
                         "20 test cases — dhammaan PASS ✓",
                         "Integration: Frontend→Backend→AI",
                         "Security: JWT, bcrypt, RBAC",
                         "Ka dhakhso badan kala saarid gacanta",
                     ])

    slide_content(prs, 10, "Gabagabo & Mustaqbal (Conclusion)", [
        "BAREAI wuxuu xallinayaa kala saarid gacanta iyadoo la adeegsanayo NLP + ML.",
        "Decision-support tool — bini'aadamku go'aanka ugu dambeeya wuu qaataa.",
        "Mustaqbal: Multi-class crime types, dataset Soomaali ballaaran.",
        "Afri-BERTa fine-tuning, mobile app, Netlify deploy.",
    ])

    slide_thank_you(prs)

    prs.save(str(OUT_FILE))
    return OUT_FILE


if __name__ == "__main__":
    path = build_pptx()
    print(f"Saved: {path}")
